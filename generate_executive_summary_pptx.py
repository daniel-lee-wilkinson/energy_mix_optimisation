from __future__ import annotations

from datetime import date
from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
FIG_DIR = ROOT / "figures"
OUT_FILE = ROOT / "Executive_Summary_DAC_Energy_Pipeline.pptx"
SENSITIVITY_FILE = ROOT / "output_data" / "demand_sensitivity_results.csv"
POTENTIALS_FILE = ROOT / "output_data" / "australian_energy_data_potentials.csv"


def set_title_style(text_frame, font_size: int = 32) -> None:
    p = text_frame.paragraphs[0]
    p.font.size = Pt(font_size)
    p.font.bold = True


def set_subtitle_style(text_frame, font_size: int = 16) -> None:
    p = text_frame.paragraphs[0]
    p.font.size = Pt(font_size)


def add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    set_title_style(slide.shapes.title.text_frame)

    subtitle_shape = slide.placeholders[1]
    subtitle_shape.text = subtitle
    set_subtitle_style(subtitle_shape.text_frame)


def add_bullet_slide(prs: Presentation, title: str, bullets: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title

    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    first = True
    for line in bullets:
        if first:
            p = body.paragraphs[0]
            first = False
        else:
            p = body.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(20)


def add_figure_slide(
    prs: Presentation,
    title: str,
    image_path: Path,
    caption: str,
    left: float = 0.8,
    top: float = 1.4,
    width: float = 11.7,
    caption_top: float = 6.8,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title

    if image_path.exists():
        slide.shapes.add_picture(
            str(image_path),
            Inches(left),
            Inches(top),
            width=Inches(width),
        )
        caption_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(caption_top), Inches(11.5), Inches(0.6)
        )
        caption_tf = caption_box.text_frame
        caption_tf.text = caption
        caption_tf.paragraphs[0].font.size = Pt(14)
    else:
        fallback = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.5), Inches(1.5))
        fallback.text_frame.text = f"Figure not found: {image_path.name}"


def build_results_bullets() -> list[str]:
    if not SENSITIVITY_FILE.exists():
        return [
            "Demand sensitivity file not found; rerun run_demand_sensitivity.py to refresh metrics.",
            "Latest optimisation outputs still available in output_data/optimal_supply_profile.csv.",
        ]

    df = pd.read_csv(SENSITIVITY_FILE)
    if df.empty:
        return [
            "Demand sensitivity file is empty.",
            "Run run_demand_sensitivity.py and complete prompts to populate this table.",
        ]

    demand_col = "annual_demand_mwh"

    net_col_candidates = [
        "net_gwp_kgco2e_per_kwh",
        "net_gwp_kg_per_kwh",
        "minimum_gwp_kgco2_per_kwh",
        "minimum_gwp_kgco2e_per_kwh",
    ]
    net_col = next((col for col in net_col_candidates if col in df.columns), None)
    if net_col is None:
        return [
            "Could not find a net/GWP metric column in demand_sensitivity_results.csv.",
            f"Available columns: {', '.join(df.columns)}",
        ]

    grid_col = "grid_share_pct"

    low = df.loc[df[demand_col].idxmin()]
    high = df.loc[df[demand_col].idxmax()]
    best = df.loc[df[net_col].idxmin()]

    bullets = [
        f"Demand sweep: {df[demand_col].min():,.0f} to {df[demand_col].max():,.0f} MWh/year ({len(df)} demand points).",
        f"Best net GWP in sweep: {best[net_col]:.3f} kg CO2e/kWh at {best[demand_col]:,.0f} MWh/year.",
    ]

    if grid_col in df.columns:
        bullets.append(
            f"Grid share trend: {low[grid_col]:.1f}% at low demand -> {high[grid_col]:.1f}% at high demand."
        )

    if "pv_share_pct" in df.columns and "wind_share_pct" in df.columns:
        bullets.append(
            f"At best net GWP point: PV {best['pv_share_pct']:.1f}%, Wind {best['wind_share_pct']:.1f}%, Grid {best.get('grid_share_pct', 0):.1f}%."
        )

    return bullets


def build_land_limit_bullets() -> list[str]:
    """Summarize how land constraints shape the optimal PV/Wind outcome."""
    land_limit_km2 = 1000.0
    pv_land_per_mw = 0.02
    wind_land_per_mw = 0.26

    max_pv_mw = land_limit_km2 / pv_land_per_mw
    max_wind_mw = land_limit_km2 / wind_land_per_mw

    bullets = [
        f"Hard constraint in example case: total land available = {land_limit_km2:,.0f} km².",
        f"Land intensity: PV {pv_land_per_mw:.2f} km²/MW vs Wind {wind_land_per_mw:.2f} km²/MW (~{wind_land_per_mw / pv_land_per_mw:.0f}x larger for wind).",
        f"Capacity implied by land limit: PV up to {max_pv_mw:,.0f} MW vs Wind up to {max_wind_mw:,.0f} MW.",
        "Optimization objective minimizes gross GWP of delivered electricity while meeting demand hourly.",
    ]

    if POTENTIALS_FILE.exists():
        df = pd.read_csv(POTENTIALS_FILE)
        if not df.empty and {"pv_potential_per_mw", "wind_potential_per_mw"}.issubset(df.columns):
            pv_per_mw_mwh = df["pv_potential_per_mw"].sum() / 1000
            wind_per_mw_mwh = df["wind_potential_per_mw"].sum() / 1000
            pv_per_km2_mwh = pv_per_mw_mwh / pv_land_per_mw if pv_land_per_mw > 0 else 0
            wind_per_km2_mwh = wind_per_mw_mwh / wind_land_per_mw if wind_land_per_mw > 0 else 0
            if wind_per_km2_mwh > 0:
                ratio = pv_per_km2_mwh / wind_per_km2_mwh
                bullets.append(
                    f"Site data signal: PV delivers ~{ratio:.1f}x more annual energy per km² than wind in current inputs."
                )

    bullets.append(
        "Resulting behavior: wind is often excluded because allocating scarce land to PV displaces more high-GWP grid energy."
    )
    return bullets


def main() -> None:
    prs = Presentation()

    add_title_slide(
        prs,
        "Executive Summary: Renewable Energy Mix for DAC Operations",
        f"Energy Mix Optimisation Pipeline | Generated {date.today().isoformat()}",
    )

    add_bullet_slide(
        prs,
        "Why This Matters for DAC",
        [
            "Electricity carbon intensity can materially reduce net CO2 removal performance.",
            "This workflow identifies the lowest-GWP supply mix under local resource and land constraints.",
            "Outputs are decision-ready: capacity targets, expected grid reliance, and net GWP impact.",
        ],
    )

    add_bullet_slide(
        prs,
        "Pipeline Overview",
        [
            "Ingest hourly NASA POWER weather data (solar irradiation, wind speed, temperature).",
            "Convert weather into hourly PV and wind generation potential per MW installed.",
            "Simulate hourly battery dispatch and residual grid imports to meet demand.",
            "Grid-search PV/Wind capacities to minimize lifecycle GWP per kWh under land limits.",
            "Publish data tables and figures for operational and strategic review.",
        ],
    )

    add_bullet_slide(
        prs,
        "Core Assumptions",
        [
            "PV: 15% conversion efficiency, temperature derating above 25 C, 0.07 kg CO2e/kWh.",
            "Wind: standard power curve (3/12/25 m/s), 0.011 kg CO2e/kWh.",
            "Land use factors: PV 0.02 km²/MW; Wind 0.26 km²/MW.",
            "Battery: 120 MWh, 90 MW, 90% discharge efficiency, 20% minimum state of charge.",
            "Grid: always available to close deficits at 0.6 kg CO2e/kWh.",
            "Demand profile: hourly scaling to annual target with heat-driven uplift above 30 C.",
        ],
    )

    add_bullet_slide(prs, "Land Limitation: Why Wind Can Be Zero", build_land_limit_bullets())

    add_bullet_slide(
        prs,
        "Optimisation Logic and KPIs",
        [
            "Primary objective: minimize gross average GWP of delivered electricity (kg CO2e/kWh).",
            "Net reporting: subtract configured annual atmospheric removal credit (default 5000 tCO2e/year).",
            "Key outputs: optimal PV/Wind MW, source shares, annual emissions, net GWP, curtailment.",
            "Sensitivity mode: evaluate how source shares and net GWP shift with demand level.",
        ],
    )

    add_bullet_slide(prs, "Current Results Snapshot", build_results_bullets())

    add_figure_slide(
        prs,
        "Demand Sensitivity: Energy Mix",
        FIG_DIR / "demand_sensitivity_energy_mix.png",
        "Source share changes as annual demand increases.",
    )

    add_figure_slide(
        prs,
        "Demand Sensitivity: Capacity Requirements",
        FIG_DIR / "demand_sensitivity_capacities.png",
        "Installed PV/Wind capacity and resulting net GWP across demand points.",
    )

    add_figure_slide(
        prs,
        "Operational Profile (Example Optimal System)",
        FIG_DIR / "optimal_system_profiles.png",
        "Hourly patterns for renewable generation, battery use, and demand coverage.",
        top=1.3,
        width=11.8,
        caption_top=6.9,
    )

    add_bullet_slide(
        prs,
        "Business Value to a DAC Company",
        [
            "Improves confidence in net-removal claims by tying energy sourcing to emissions outcomes.",
            "Supports land and infrastructure planning with transparent capacity-to-performance trade-offs.",
            "Enables procurement strategy for PPAs and firming based on modeled grid dependence.",
            "Provides a repeatable baseline for techno-economic and financing due diligence.",
        ],
    )

    add_bullet_slide(
        prs,
        "Recommended Next Steps",
        [
            "Co-optimise battery size (MWh/MW) with PV/Wind capacity to reduce grid reliance further.",
            "Add dynamic grid emissions and electricity pricing to evaluate cost-carbon frontiers.",
            "Run uncertainty ranges for weather years and component performance degradation.",
            "Integrate this model with plant-level DAC operations for dispatch-aware optimization.",
        ],
    )

    prs.save(OUT_FILE)
    print(f"Saved: {OUT_FILE}")


if __name__ == "__main__":
    main()
